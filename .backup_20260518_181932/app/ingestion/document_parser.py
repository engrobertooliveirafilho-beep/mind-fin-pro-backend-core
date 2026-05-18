import io
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

class DocumentParser:

    def parse_pdf(self, path):
        text = []
        reader = PdfReader(path)
        for page in reader.pages:
            text.append(page.extract_text() or "")
        return "\n".join(text)

    def parse_docx(self, path):
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])

    def parse_pptx(self, path):
        prs = Presentation(path)
        output = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    output.append(shape.text)

        return "\n".join(output)

    def parse_xlsx(self, path):
        xls = pd.ExcelFile(path)
        output = []

        for sheet in xls.sheet_names:
            df = xls.parse(sheet)
            output.append(f"PLANILHA: {sheet}")
            output.append(df.head(20).to_string())

        return "\n".join(output)

    def parse_txt(self, path):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
